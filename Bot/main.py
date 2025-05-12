from googleapiclient.discovery import build
import os
import json
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
import google.generativeai as genai
from imapclient import IMAPClient
import pyzmail
import requests
from bs4 import BeautifulSoup
import re
import tkinter as tk
from tkinter import filedialog
from docx import Document
import PyPDF2
from pptx import Presentation
import pandas as pd
from tkinter import messagebox
import io
import pytesseract
from PIL import Image
from graphviz import Digraph
from rich.console import Console
from rich.markdown import Markdown
from rich.panel import Panel
from rich.rule import Rule
from rich.table import Table
from rich.columns import Columns
from rich.style import Style
from rich.text import Text
from rich.box import SIMPLE
from rich.tree import Tree
from pathlib import Path
import os
import requests
from io import BytesIO
from PIL import Image
import base64
import shutil

console = Console()

# Define custom styles for consistent and vibrant output
STYLE_PROMPT = Style(color="cyan", bold=True)
STYLE_BOT = Style(color="bright_yellow")
STYLE_ERROR = Style(color="red", bold=True)
STYLE_WARNING = Style(color="yellow", italic=True)
STYLE_INFO = Style(color="bright_blue")
STYLE_COMMAND = Style(color="magenta", bold=True)
STYLE_SUCCESS = Style(color="green", bold=True)

# ---------- History utilities ----------
HISTORY_FILE = Path(os.getenv("HISTORY_FILE", Path(__file__).parent / "chat_history.json"))

def load_history():
    """Loads chat history from a JSON file."""
    raw = []
    gemini = []
    if HISTORY_FILE.exists():
        try:
            raw = json.loads(HISTORY_FILE.read_text(encoding="utf-8"))
            for entry in raw:
                if "user" in entry and "bot" in entry:
                    gemini.append({"role": "user",  "parts": [entry["user"]]})
                    gemini.append({"role": "model", "parts": [entry["bot"]]})
                else:
                    console.print(Panel(
                        "Skipping malformed history entry.",
                        style=STYLE_WARNING,
                        border_style=STYLE_WARNING,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
        except json.JSONDecodeError:
            console.print(Panel(
                "Could not parse history file (invalid JSON); starting fresh.",
                style=STYLE_ERROR,
                border_style=STYLE_ERROR,
                box=SIMPLE,
                padding=(0, 1)
            ))
            raw, gemini = [], []
        except Exception as e:
            console.print(Panel(
                f"Could not load history file at {HISTORY_FILE}: {e}; starting fresh.",
                style=STYLE_ERROR,
                border_style=STYLE_ERROR,
                box=SIMPLE,
                padding=(0, 1)
            ))
            raw, gemini = [], []
    else:
        console.print(Panel(
            f"History file not found at {HISTORY_FILE}; starting fresh.",
            style=STYLE_WARNING,
            border_style=STYLE_WARNING,
            box=SIMPLE,
            padding=(0, 1)
        ))
    return raw, gemini

def save_history(raw_history):
    """Saves chat history to a JSON file."""
    try:
        HISTORY_FILE.parent.mkdir(parents=True, exist_ok=True)
        HISTORY_FILE.write_text(
            json.dumps(raw_history, indent=2, ensure_ascii=False), encoding="utf-8"
        )
    except Exception as e:
        console.print(Panel(
            f"Could not save history to {HISTORY_FILE}: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))

# ---------- Browse Function ----------
def browse(query, url):
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        for script in soup(["script", "style"]):
            script.decompose()
        content_tags = soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'article', 'div'])
        text_parts = []
        for tag in content_tags:
            text = tag.get_text(strip=True)
            if text and len(text) > 20:
                text_parts.append(text)
        full_text = " ".join(text_parts)
        full_text = re.sub(r'\s+', ' ', full_text).strip()
        if not full_text:
            return "Fetched page, but no readable content found."
        return full_text[:5000] + "..." if len(full_text) > 5000 else full_text
    except requests.exceptions.RequestException as e:
        return f"Error fetching page: {str(e)}"
    except Exception as e:
        return f"Error processing page content: {str(e)}"

# ---------- Google Search Function ----------
my_api_key = "YOUR_GOOGLE_API_KEY"
my_cse_id = "YOUR_CSE_ID"

def google_search(search_term, api_key, cse_id, **kwargs):
    service = build("customsearch", "v1", developerKey=api_key)
    try:
        res = service.cse().list(q=search_term, cx=cse_id, **kwargs).execute()
        items = res.get('items', [])
        enriched_results = []
        for item in items:
            result = {
                'title': item.get('title', 'No Title'),
                'snippet': item.get('snippet', 'No Snippet'),
                'url': item.get('link', 'No URL'),
                'content': 'Could not fetch page content.',
                'content_status': 'Not attempted (no valid URL)'
            }
            url = result['url']
            if url and url.startswith(('http://', 'https://')):
                console.print(Panel(
                    f"Fetching content from: {url}",
                    style=STYLE_INFO,
                    border_style=STYLE_INFO,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                content = browse(query=search_term, url=url)
                if content and content.strip() and not content.startswith('Error'):
                    result['content'] = content
                    console.print(Panel(
                        f"✅ Successfully fetched content from {url}",
                        style=STYLE_SUCCESS,
                        border_style=STYLE_SUCCESS,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    result['content_status'] = 'Success'
                else:
                    result['content'] = content if content else 'Fetched page, but no readable content found.'
                    console.print(Panel(
                        f"⚠️ {result['content']} at {url}",
                        style=STYLE_WARNING,
                        border_style=STYLE_WARNING,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    result['content_status'] = 'No readable content' if not content.startswith('Error') else 'Fetch Error'
            else:
                console.print(Panel(
                    f"Skipping invalid or missing URL: {url}",
                    style=STYLE_WARNING,
                    border_style=STYLE_WARNING,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
            enriched_results.append(result)
        return enriched_results
    except Exception as e:
        console.print(Panel(
            f"Error during Google Search API call: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return []

# ---------- Constants ----------
EMAIL = os.getenv("EMAIL", 'YOUR_EMAIL')
PASSWORD = os.getenv("PASSWORD", 'YOUR_PASSWORD')
IMAP_SERVER = os.getenv("IMAP_SERVER", 'YOUR_IMAP_SERVER')
NEWS_API_KEY = os.getenv("NEWS_API_KEY", 'YOUR_NEWS_API_KEY')

# --- Stock code ---
def run_get_stocks():
    ALPHA_VANTAGE_API_KEY = os.getenv("ALPHA_VANTAGE_API_KEY")
    if not ALPHA_VANTAGE_API_KEY:
        console.print(Panel(
            "Error: ALPHA_VANTAGE_API_KEY not found in .env file.\n\n"
            "To fetch stock prices, please:\n"
            "1. Ensure the .env file is in the script's directory or correctly referenced.\n"
            "2. Add the key like this: ALPHA_VANTAGE_API_KEY=your_api_key_here\n"
            "3. Verify the .env file has no extra spaces or comments before the key.\n"
            "4. Obtain a free API key from https://www.alphavantage.co/support/#api-key if needed.",
            title="Stock Fetching Error",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(1, 2)
        ))
        console.print(Panel(
            f"Debug: Checked .env at {Path(__file__).parent / '.env'}, ALPHA_VANTAGE_API_KEY={os.getenv('ALPHA_VANTAGE_API_KEY')}",
            style=STYLE_INFO,
            border_style=STYLE_INFO,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return
    STOCKS = ["AAPL", "GOOGL", "AMZN", "MSFT", "BTCUSD"]

    console.print(Rule("[bold bright_yellow]Stock Prices[/]", style="bright_yellow"))
    console.print(Panel(
        "Fetching stock prices...",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))

    stock_data = []
    for symbol in STOCKS:
        try:
            resp = requests.get('https://www.alphavantage.co/query', params={
                'function': 'TIME_SERIES_DAILY_ADJUSTED', 'symbol': symbol, 'apikey': ALPHA_VANTAGE_API_KEY
            })
            resp.raise_for_status()
            js = resp.json()
            ts = js.get('Time Series (Daily Adjusted)')
            if not ts:
                console.print(Panel(
                    f"No daily time series data for {symbol}. Response: {js.get('Note', 'No data.')}",
                    style=STYLE_WARNING,
                    border_style=STYLE_WARNING,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                stock_data.append({'symbol': symbol, 'last_price': 'N/A', 'price_change': 'N/A'})
                continue
            dates = sorted(ts.keys(), reverse=True)
            if not dates:
                console.print(Panel(
                    f"No dates found in data for {symbol}.",
                    style=STYLE_WARNING,
                    border_style=STYLE_WARNING,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                stock_data.append({'symbol': symbol, 'last_price': 'N/A', 'price_change': 'N/A'})
                continue
            last_price_str = ts[dates[0]]['4. close']
            last = float(last_price_str)
            if len(dates) > 29:
                prev = float(ts[dates[29]]['4. close'])
                price_change = (last - prev) / prev * 100
            else:
                price_change = 'N/A'
                console.print(Panel(
                    f"Less than 30 days of data for {symbol}. Cannot calculate 30-day change.",
                    style=STYLE_WARNING,
                    border_style=STYLE_WARNING,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
            stock_data.append({'symbol': symbol, 'last_price': last, 'price_change': price_change})
        except requests.exceptions.RequestException as e:
            console.print(Panel(
                f"Error fetching data for {symbol}: {e}",
                style=STYLE_ERROR,
                border_style=STYLE_ERROR,
                box=SIMPLE,
                padding=(0, 1)
            ))
            stock_data.append({'symbol': symbol, 'last_price': 'Error', 'price_change': 'Error'})
        except Exception as e:
            console.print(Panel(
                f"Error processing data for {symbol}: {e}",
                style=STYLE_ERROR,
                border_style=STYLE_ERROR,
                box=SIMPLE,
                padding=(0, 1)
            ))
            stock_data.append({'symbol': symbol, 'last_price': 'Error', 'price_change': 'Error'})
    
    if stock_data:
        table = Table(
            title="Current Stock Prices & Changes (vs ~30 days ago)",
            style="green",
            box=SIMPLE,
            title_style="bold bright_green",
            header_style="bold bright_cyan"
        )
        table.add_column("Symbol", style="bold bright_white")
        table.add_column("Last Price (USD)", justify="right", style="white")
        table.add_column("Change (%)", justify="right", style="white")
        for i in stock_data:
            change_str = f"{i['price_change']:.2f}%" if isinstance(i['price_change'], (int, float)) else str(i['price_change'])
            change_color = "green" if isinstance(i['price_change'], (int, float)) and i['price_change'] >= 0 else "red" if isinstance(i['price_change'], (int, float)) else "yellow"
            last_price_str = f"{i['last_price']:.2f}" if isinstance(i['last_price'], (int, float)) else str(i['last_price'])
            table.add_row(i['symbol'], last_price_str, f"[{change_color}]{change_str}[/]")
        console.print(Panel(
            table,
            border_style="bright_green",
            box=SIMPLE,
            padding=(0, 1)
        ))
    else:
        console.print(Panel(
            "No stock data retrieved.",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
    console.print(Rule(style="bright_yellow"))

# --- File extraction ---
def extract_text_from_file():
    root = tk.Tk()
    root.withdraw()
    fp = filedialog.askopenfilename(title="Select file", filetypes=[
        ("All supported files", "*.txt *.pdf *.docx *.pptx *.xlsx *.py *.png *.jpg *.jpeg"),
        ("Text files", "*.txt"),
        ("PDF files", "*.pdf"),
        ("Word documents", "*.docx"),
        ("PowerPoint presentations", "*.pptx"),
        ("Excel spreadsheets", "*.xlsx"),
        ("Python scripts", "*.py"),
        ("Image files", "*.png *.jpg *.jpeg")
    ])
    if not fp:
        console.print(Panel(
            "No file selected.",
            style=STYLE_WARNING,
            border_style=STYLE_WARNING,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return None
    console.print(Panel(
        f"Processing file: [bold]{Path(fp).name}[/bold]…",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))
    ext = os.path.splitext(fp)[1].lower()
    text = ""
    try:
        if ext in ('.txt', '.py'):
            with open(fp, 'r', encoding='utf-8') as f:
                text = f.read()
        elif ext == '.pdf':
            with open(fp, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                if reader.is_encrypted:
                    console.print(Panel(
                        "Error: PDF file is encrypted. Cannot read.",
                        style=STYLE_ERROR,
                        border_style=STYLE_ERROR,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    return None
                page_texts = []
                for i, p in enumerate(reader.pages):
                    try:
                        page_texts.append(p.extract_text() or '')
                    except Exception as page_e:
                        console.print(Panel(
                            f"Warning: Could not extract text from page {i+1}: {page_e}",
                            style=STYLE_WARNING,
                            border_style=STYLE_WARNING,
                            box=SIMPLE,
                            padding=(0, 1)
                        ))
                        page_texts.append('')
                text = "\n".join(page_texts)
        elif ext == '.docx':
            doc = Document(fp)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == '.pptx':
            prs = Presentation(fp)
            slide_texts = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                slide_texts.append("\n".join(slide_text))
            text = "\n\n".join(slide_texts)
        elif ext == '.xlsx':
            try:
                xl = pd.ExcelFile(fp)
                sheet_texts = []
                for sheet_name in xl.sheet_names:
                    df = pd.read_excel(fp, sheet_name=sheet_name)
                    sheet_texts.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}")
                text = "\n\n".join(sheet_texts)
            except Exception as excel_e:
                console.print(Panel(
                    f"Error reading Excel file: {excel_e}",
                    style=STYLE_ERROR,
                    border_style=STYLE_ERROR,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                return None
        elif ext in ('.png', '.jpg', '.jpeg'):
            try:
                api_key = os.environ.get("API_KEY")
                if not api_key:
                    console.print(Panel(
                        "Error: API_KEY not found for Gemini image analysis.",
                        style=STYLE_ERROR,
                        border_style=STYLE_ERROR,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    return None
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel('gemini-1.5-flash')
                img = Image.open(fp)
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format=img.format)
                img_bytes = img_byte_arr.getvalue()
                prompt = (
                    "Analyze this image in detail. Describe the objects, text, scenes, or any notable features. "
                    "If text is present, include it in your description. Provide a comprehensive summary."
                )
                response = model.generate_content([prompt, {"mime_type": f"image/{ext[1:]}", "data": img_bytes}])
                text = response.text
                console.print(Panel(
                    "✅ Image analyzed successfully.",
                    style=STYLE_SUCCESS,
                    border_style=STYLE_SUCCESS,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
            except Exception as gemini_e:
                console.print(Panel(
                    f"Error analyzing image with Gemini: {gemini_e}",
                    style=STYLE_ERROR,
                    border_style=STYLE_ERROR,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                return None
        else:
            console.print(Panel(
                f"Unsupported file type: {ext}",
                style=STYLE_ERROR,
                border_style=STYLE_ERROR,
                box=SIMPLE,
                padding=(0, 1)
            ))
            return None
        if text.strip():
            console.print(Panel(
                "✅ Processing complete.",
                style=STYLE_SUCCESS,
                border_style=STYLE_SUCCESS,
                box=SIMPLE,
                padding=(0, 1)
            ))
        else:
            console.print(Panel(
                "⚠️ File processed, but no readable content was found.",
                style=STYLE_WARNING,
                border_style=STYLE_WARNING,
                box=SIMPLE,
                padding=(0, 1)
            ))
            return None
    except FileNotFoundError:
        console.print(Panel(
            f"Error: File not found at {fp}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return None
    except Exception as e:
        console.print(Panel(
            f"Error processing file: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return None
    return text 

# --- Email & News ---
def fetch_emails():
    summs = []
    console.print(Rule("[bold bright_blue]Fetching Emails[/]", style="bright_blue"))
    console.print(Panel(
        "Connecting to email server…",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))
    try:
        with IMAPClient(IMAP_SERVER) as c:
            c.login(EMAIL, PASSWORD)
            c.select_folder('INBOX', readonly=True)
            msgs = c.search(['ALL'])
            recent_msgs = msgs[-10:] if len(msgs) > 10 else msgs
            if not recent_msgs:
                console.print(Panel(
                    "No emails found in INBOX.",
                    style=STYLE_WARNING,
                    border_style=STYLE_WARNING,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                return ["No emails found."]
            console.print(Panel(
                f"Fetching data for {len(recent_msgs)} recent emails…",
                style=STYLE_INFO,
                border_style=STYLE_INFO,
                box=SIMPLE,
                padding=(0, 1)
            ))
            fetch_response = c.fetch(recent_msgs, ['BODY.PEEK[]', 'ENVELOPE'])
            for mid in recent_msgs:
                msg_data = fetch_response.get(mid)
                if not msg_data:
                    console.print(Panel(
                        f"Could not fetch data for message ID {mid}.",
                        style=STYLE_WARNING,
                        border_style=STYLE_WARNING,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    continue
                envelope = msg_data.get(b'ENVELOPE')
                subj = 'No Subject'
                frm = 'Unknown Sender'
                if envelope:
                    if envelope.subject:
                        try:
                            subj = envelope.subject.decode()
                        except Exception:
                            subj = 'Undecodable Subject'
                    if envelope.from_:
                        frm_address = envelope.from_[0]
                        try:
                            name = frm_address.name.decode() if frm_address.name else ''
                            mailbox = frm_address.mailbox.decode() if frm_address.mailbox else ''
                            host = frm_address.host.decode() if frm_address.host else ''
                            if mailbox and host:
                                frm = f"{name} <{mailbox}@{host}>".strip()
                            elif mailbox:
                                frm = mailbox
                            elif name:
                                frm = name
                        except Exception:
                            frm = 'Undecodable Sender Info'
                body = 'No readable body.'
                raw_body = msg_data.get(b'BODY[]')
                if raw_body:
                    try:
                        m = pyzmail.PyzMessage.factory(raw_body)
                        part = m.text_part or m.html_part
                        if part:
                            try:
                                body = part.get_payload().decode(part.charset if part.charset else 'utf-8', errors='replace')
                            except Exception as decode_error:
                                body = f"Could not decode email body: {decode_error}"
                                console.print(Panel(
                                    f"Warning: Could not decode email body for message ID {mid}.",
                                    style=STYLE_WARNING,
                                    border_style=STYLE_WARNING,
                                    box=SIMPLE,
                                    padding=(0, 1)
                                ))
                            if subj == 'No Subject' and m.get_subject():
                                subj = m.get_subject()
                            if frm == 'Unknown Sender' and m.get_addresses('from'):
                                frm = m.get_addresses('from')[0][1]
                    except Exception as pyzmail_error:
                        console.print(Panel(
                            f"Warning: Could not parse email with pyzmail for message ID {mid}: {pyzmail_error}",
                            style=STYLE_WARNING,
                            border_style=STYLE_WARNING,
                            box=SIMPLE,
                            padding=(0, 1)
                        ))
                body_preview = body[:500] + "..." if len(body) > 500 else body
                console.print(Panel(
                    f"From: {frm}\nSubject: {subj}\nBody Preview: {body_preview}",
                    style=STYLE_INFO,
                    border_style=STYLE_INFO,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                summs.append(f"From: {frm}\nSubject: {subj}\nBody: {body}")
            console.print(Panel(
                "✅ Email fetching complete.",
                style=STYLE_SUCCESS,
                border_style=STYLE_SUCCESS,
                box=SIMPLE,
                padding=(0, 1)
            ))
        return summs
    except Exception as e:
        console.print(Panel(
            f"Error fetching emails: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return [f"Error fetching emails: {e}"]
    finally:
        console.print(Rule(style="bright_blue"))

def fetch_news():
    console.print(Rule("[bold magenta]Fetching News[/]", style="magenta"))
    console.print(Panel(
        "Fetching top headlines…",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))
    try:
        r = requests.get(f'https://newsapi.org/v2/top-headlines?country=us&pageSize=100&apiKey={NEWS_API_KEY}')
        r.raise_for_status()
        articles = r.json().get('articles', [])
        enriched_articles = []
        count = 0
        for article in articles:
            if article and isinstance(article, dict) and article.get('title') and article.get('description'):
                enriched_article = {
                    'title': article.get('title', 'No Title'),
                    'description': article.get('description', 'No Description'),
                    'url': article.get('url', 'No URL'),
                    'content': 'Could not fetch article content.',
                    'content_status': 'Not attempted (no valid URL)'
                }
                url = enriched_article['url']
                if url and url.startswith(('http://', 'https://')):
                    console.print(Panel(
                        f"Fetching content from: {url}",
                        style=STYLE_INFO,
                        border_style=STYLE_INFO,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    content = browse(query="news", url=url)
                    if content and content.strip() and not content.startswith('Error'):
                        enriched_article['content'] = content
                        console.print(Panel(
                            f"✅ Fetched content from {url}",
                            style=STYLE_SUCCESS,
                            border_style=STYLE_SUCCESS,
                            box=SIMPLE,
                            padding=(0, 1)
                        ))
                        enriched_article['content_status'] = 'Success'
                    else:
                        enriched_article['content'] = content if content else 'Fetched page, but no readable content found.'
                        console.print(Panel(
                            f"⚠️ {enriched_article['content']} at {url}",
                            style=STYLE_WARNING,
                            border_style=STYLE_WARNING,
                            box=SIMPLE,
                            padding=(0, 1)
                        ))
                        enriched_article['content_status'] = 'No readable content' if not content.startswith('Error') else 'Fetch Error'
                else:
                    console.print(Panel(
                        f"Skipping invalid or missing URL: {url}",
                        style=STYLE_WARNING,
                        border_style=STYLE_WARNING,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                formatted_article = (
                    f"Title: {enriched_article['title']}\n"
                    f"Description: {enriched_article['description']}\n"
                    f"URL: {enriched_article['url']}\n"
                    f"Content Fetch Status: {enriched_article['content_status']}\n"
                    f"Article Content Preview:\n{enriched_article['content']}"
                )
                enriched_articles.append(formatted_article)
                count += 1
                if count >= 200:
                    break
        if not enriched_articles:
            console.print(Panel(
                "No news articles found with both title and description.",
                style=STYLE_WARNING,
                border_style=STYLE_WARNING,
                box=SIMPLE,
                padding=(0, 1)
            ))
            return ["No news articles found."]
        console.print(Panel(
            f"✅ Fetched {len(enriched_articles)} relevant news articles.",
            style=STYLE_SUCCESS,
            border_style=STYLE_SUCCESS,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return enriched_articles
    except requests.exceptions.RequestException as e:
        console.print(Panel(
            f"Error fetching news: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return [f"Error fetching news: {e}"]
    except Exception as e:
        console.print(Panel(
            f"Error processing news data: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return [f"Error processing news data: {e}"]
    finally:
        console.print(Rule(style="magenta"))

# --- Wikipedia fetch ---
def fetch_wikipedia(topic: str) -> str:
    console.print(Rule("[bold bright_green]Fetching Wikipedia Info[/]", style="bright_green"))
    console.print(Panel(
        f"Fetching Wikipedia summary for '[bold]{topic}[/bold]'…",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))
    url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{requests.utils.quote(topic.replace(' ','_'))}"
    try:
        r = requests.get(url)
        r.raise_for_status()
        data = r.json()
        extract = data.get('extract')
        if extract:
            console.print(Panel(
                "✅ Wikipedia summary retrieved successfully.",
                style=STYLE_SUCCESS,
                border_style=STYLE_SUCCESS,
                box=SIMPLE,
                padding=(0, 1)
            ))
            return extract
        else:
            title = data.get('title', 'N/A')
            description = data.get('description', 'N/A')
            console.print(Panel(
                f"No direct summary found for '[bold]{topic}[/bold]'. Found: [bold]{title}[/bold] ({description})",
                style=STYLE_WARNING,
                border_style=STYLE_WARNING,
                box=SIMPLE,
                padding=(0, 1)
            ))
            if title and description != 'N/A':
                return f"Wikipedia found something related to '{topic}', possibly: {title} ({description}). No direct summary extract available."
            elif title != 'N/A':
                return f"Wikipedia found a page titled '{title}' for '{topic}', but no summary extract."
            else:
                return 'No summary found for this topic on Wikipedia.'
    except requests.exceptions.RequestException as e:
        console.print(Panel(
            f"Error fetching Wikipedia article: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return f'Failed to retrieve Wikipedia article from API: {e}'
    except Exception as e:
        console.print(Panel(
            f"Error processing Wikipedia data: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return f'Failed to process Wikipedia data: {e}'
    finally:
        console.print(Rule(style="bright_green"))

# --- Mindmap Generator ---
def generate_mindmap(model, prompt, raw_hist, file_content=None):
    """
    Generate a detailed graphical mindmap with a balanced radial layout, lighter colors, and classic style,
    and automatically open the resulting image. Adjusted for new script location at C:\Bot.
    Args:
        model: Gemini GenerativeModel instance.
        prompt: User-provided prompt for the mindmap.
        raw_hist: List of chat history entries.
        file_content: Optional content extracted from a file.
    Returns:
        Path to the generated mindmap image (PNG) or error message.
    """
    # Specify Graphviz bin path
    graphviz_bin_path = r"C:\Program Files\Graphviz\bin"
    
    # Verify dot.exe exists
    dot_exe_path = os.path.join(graphviz_bin_path, "dot.exe")
    if not os.path.exists(dot_exe_path):
        console.print(Panel(
            f"Error: 'dot.exe' not found at {dot_exe_path}. Please ensure Graphviz is installed correctly.",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return "Failed to generate mindmap: Graphviz 'dot.exe' not found."

    # Temporarily add Graphviz to PATH
    os.environ["PATH"] = graphviz_bin_path + os.pathsep + os.environ["PATH"]
    console.print(Panel(
        f"Debug: Temporarily added Graphviz bin to PATH: {graphviz_bin_path}",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))

    # Verify dot.exe is found
    dot_path = shutil.which('dot')
    console.print(Panel(
        f"Debug: Graphviz 'dot' executable found at: {dot_path if dot_path else 'Not found'}",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))
    if not dot_path:
        return "Failed to render mindmap: Graphviz 'dot' executable not found in PATH."

    console.print(Panel(
        f"Generating graphical mindmap for prompt: [bold]{prompt[:50]}{'...' if len(prompt) > 50 else ''}[/bold]",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))

    # Correct common typos
    if prompt.lower() == "imerialism":
        prompt = "imperialism"
        console.print(Panel(
            f"Corrected typo: 'imerialism' to 'imperialism'",
            style=STYLE_WARNING,
            border_style=STYLE_WARNING,
            box=SIMPLE,
            padding=(0, 1)
        ))

    # Summarize recent chat history
    history_summary = ""
    if raw_hist:
        recent_hist = raw_hist[-5:]
        history_lines = []
        for entry in recent_hist:
            history_lines.append(f"User ({entry['timestamp']}): {entry['user']}")
            history_lines.append(f"Bot: {entry['bot']}\n")
        history_summary = "\n".join(history_lines)
        if len(history_summary) > 1000:
            history_summary = history_summary[:1000] + "..."

    # Gemini prompt (unchanged)
    gemini_prompt = (
        "You are an AI assistant that generates highly detailed, structured mindmaps. "
        f"Create a comprehensive mindmap based on the following user prompt: '{prompt}'. "
    )
    if history_summary:
        gemini_prompt += (
            f"\nUse the following chat history as context:\n```\n{history_summary}\n```"
        )
    if file_content:
        gemini_prompt += (
            f"\nUse the following file content as context:\n```\n{file_content[:2000]}\n```"
        )
    gemini_prompt += (
        "Return the result in JSON format with a central topic, 5-7 main branches, each with 3-5 sub-branches, "
        "and 1-2 sub-sub-branches per sub-branch where appropriate. "
        "Ensure the JSON is valid and follows this format:\n"
        "```json\n"
        "{\n"
        "  \"central_topic\": \"Topic\",\n"
        "  \"branches\": [\n"
        "    {\n"
        "      \"name\": \"Branch 1\",\n"
        "      \"sub_branches\": [\n"
        "        {\n"
        "          \"name\": \"Sub 1\",\n"
        "          \"sub_sub_branches\": [\"SubSub 1\", \"SubSub 2\"]\n"
        "        }\n"
        "      ]\n"
        "    }\n"
        "  ]\n"
        "}\n"
        "```"
    )

    # Fallback prompt (unchanged)
    fallback_prompt = (
        f"Generate a detailed mindmap for the prompt '{prompt}' in JSON format. "
    )
    if history_summary:
        fallback_prompt += f"\nUse the following chat history as context:\n```\n{history_summary}\n```"
    if file_content:
        fallback_prompt += f"\nUse the following file content as context:\n```\n{file_content[:1000]}\n```"
    fallback_prompt += (
        "Include a central topic, 5 main branches, each with 2-3 sub-branches, and 1 sub-sub-branch where possible. "
        "Ensure valid JSON like:\n"
        "```json\n"
        "{\n"
        "  \"central_topic\": \"Topic\",\n"
        "  \"branches\": [\n"
        "    {\"name\": \"Branch\", \"sub_branches\": [{\"name\": \"Sub\", \"sub_sub_branches\": [\"SubSub\"]}]}\n"
        "  ]\n"
        "}\n"
        "```"
    )

    mindmap_data = None
    for attempt, current_prompt in enumerate([gemini_prompt, fallback_prompt], 1):
        try:
            response = model.generate_content(current_prompt)
            raw_response = response.text.strip()
            mindmap_json = raw_response
            if mindmap_json.startswith("```json") and mindmap_json.endswith("```"):
                mindmap_json = mindmap_json[7:-3].strip()
            elif mindmap_json.startswith("```") and mindmap_json.endswith("```"):
                mindmap_json = mindmap_json[3:-3].strip()
            mindmap_data = json.loads(mindmap_json)
            if not isinstance(mindmap_data, dict) or "central_topic" not in mindmap_data or "branches" not in mindmap_data:
                raise ValueError("Invalid mindmap structure")
            console.print(Panel(
                f"✅ Valid JSON mindmap received (attempt {attempt}).",
                style=STYLE_SUCCESS,
                border_style=STYLE_SUCCESS,
                box=SIMPLE,
                padding=(0, 1)
            ))
            break
        except Exception as e:
            console.print(Panel(
                f"Error (attempt {attempt}): {e}",
                style=STYLE_ERROR,
                border_style=STYLE_ERROR,
                box=SIMPLE,
                padding=(0, 1)
            ))
            if attempt == 2:
                return "Failed to generate mindmap after multiple attempts."

    if not mindmap_data:
        return "Failed to generate a valid mindmap."

    # Create graphical mindmap with Graphviz using twopi for radial layout
    dot = Digraph(
        comment=f"Mindmap for {mindmap_data['central_topic']}",
        format='png',
        engine='twopi',  # Radial layout
        graph_attr={
            'bgcolor': '#F5F6F5',  # Light gray background
            'concentrate': 'true',  # Merge edges where possible
            'pad': '1.0',          # Padding around the graph
            'nodesep': '0.6',      # Slightly more space between nodes
            'ranksep': '1.8',      # Increased space between ranks
            'dimen': '10',         # Maximum diameter to control size
            'overlap': 'false'     # Prevent node overlap
        },
        node_attr={
            'shape': 'box',        # Rounded rectangles
            'style': 'filled',
            'fontname': 'Arial',
            'fontsize': '11',      # Slightly larger font for readability
            'margin': '0.15,0.1',  # Adjusted margins for better text fit
            'width': '1.2',        # Slightly wider nodes
            'height': '0.6'        # Slightly taller nodes
        },
        edge_attr={
            'color': '#6B7280',    # Soft gray edges
            'penwidth': '1.0',     # Thinner edges
            'arrowsize': '0.6'     # Smaller arrows
        }
    )

    # Central topic
    central_id = 'central'
    dot.node(
        central_id,
        mindmap_data['central_topic'],
        fillcolor='#BAE6FD',    # Light blue
        fontcolor='#1E3A8A',    # Dark blue text
        shape='ellipse',
        fontsize='14',
        fontname='Arial Bold'
    )

    # Color palette for branches
    branch_colors = [
        '#A7F3D0', '#FEF3C7', '#FBCFE8', '#BFDBFE',
        '#D1D5DB', '#F3E8FF', '#FECACA'
    ]

    # Add branches
    for i, branch in enumerate(mindmap_data['branches']):
        branch_id = f"branch_{i}"
        dot.node(
            branch_id,
            branch['name'],
            fillcolor=branch_colors[i % len(branch_colors)],
            fontcolor='#1F2937'
        )
        dot.edge(central_id, branch_id)

        # Add sub-branches
        for j, sub_branch in enumerate(branch.get('sub_branches', [])):
            sub_branch_id = f"sub_branch_{i}_{j}"
            dot.node(
                sub_branch_id,
                sub_branch['name'],
                fillcolor=branch_colors[i % len(branch_colors)],
                fontcolor='#1F2937',
                alpha='0.7'
            )
            dot.edge(branch_id, sub_branch_id)

            # Add sub-sub-branches
            for k, sub_sub_branch in enumerate(sub_branch.get('sub_sub_branches', [])):
                sub_sub_id = f"sub_sub_branch_{i}_{j}_{k}"
                dot.node(
                    sub_sub_id,
                    sub_sub_branch,
                    fillcolor=branch_colors[i % len(branch_colors)],
                    fontcolor='#1F2937',
                    alpha='0.5'
                )
                dot.edge(sub_branch_id, sub_sub_id)

    # Save and render the mindmap
    script_dir = Path(r"C:\Bot")  # Explicitly set to new location
    output_dir = script_dir / "mindmaps"
    try:
        output_dir.mkdir(exist_ok=True)
        console.print(Panel(
            f"Debug: Mindmap output directory set to {output_dir}",
            style=STYLE_INFO,
            border_style=STYLE_INFO,
            box=SIMPLE,
            padding=(0, 1)
        ))
    except Exception as e:
        console.print(Panel(
            f"Error creating mindmap output directory {output_dir}: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return f"Failed to create mindmap directory: {e}"

    output_file = output_dir / f"mindmap_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    try:
        dot.render(output_file, view=False, cleanup=True)
        output_path = str(output_file) + ".png"
        console.print(Panel(
            f"✅ Mindmap generated and saved to {output_path}",
            style=STYLE_SUCCESS,
            border_style=STYLE_SUCCESS,
            box=SIMPLE,
            padding=(0, 1)
        ))
        # Automatically open the mindmap image
        os.startfile(output_path)
        return output_path
    except Exception as e:
        console.print(Panel(
            f"Error rendering mindmap: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return f"Failed to render mindmap: {e}"

def generate_image(prompt):
    """Generate an image using a hypothetical image generation API."""
    console.print(Panel(
        f"Generating image for prompt: [bold]{prompt}[/bold]…",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))
    
    # Placeholder: Replace with actual image generation API (e.g., Stable Diffusion)
    api_key = os.environ.get("IMAGE_GEN_API_KEY")
    if not api_key:
        console.print(Panel(
            "Error: IMAGE_GEN_API_KEY not found in .env file.",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return None

    try:
        # Example API call (hypothetical)
        response = requests.post(
            "https://api.imagegen.example/v1/generate",
            headers={"Authorization": f"Bearer {api_key}"},
            json={"prompt": prompt, "size": "512x512", "style": "realistic"}
        )
        response.raise_for_status()
        image_data = response.json().get("image")  # Assume base64-encoded image
        image_bytes = base64.b64decode(image_data)
        image = Image.open(BytesIO(image_bytes))
        
        # Save or display the image
        output_path = Path(__file__).parent / f"generated_image_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
        image.save(output_path)
        console.print(Panel(
            f"✅ Image generated and saved to {output_path}",
            style=STYLE_SUCCESS,
            border_style=STYLE_SUCCESS,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return str(output_path)
    except Exception as e:
        console.print(Panel(
            f"Error generating image: {e}",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(0, 1)
        ))
        return None

# ---------- Main ----------
def main():
    console.print(Rule("[bold bright_cyan]SnorlaxGPT Chatbot[/]", style="bright_cyan"))
    console.print(Panel(
        "[bold bright_yellow]Hey there! SnorlaxGPT is waking up… kinda.[/bold bright_yellow]\n"
        "Ready to chat, fetch stuff, or just... be?\n"
        "[italic cyan]Don't poke me too hard, okay? Zzzz…[/italic cyan]",
        title="[bold bright_cyan]Welcome to SnorlaxGPT![/]",
        style="bright_cyan",
        border_style="bright_cyan",
        box=SIMPLE,
        padding=(1, 2)
    ))

    # Load .env and debug environment variables
    env_path = Path(__file__).parent / ".env"
    load_dotenv(dotenv_path=env_path, override=True)
    console.print(Panel(
        f"Debug: Loaded .env from {env_path}\n"
        f"ALPHA_VANTAGE_API_KEY={os.getenv('ALPHA_VANTAGE_API_KEY')}\n"
        f"HISTORY_FILE={os.getenv('HISTORY_FILE') or HISTORY_FILE}",
        style=STYLE_INFO,
        border_style=STYLE_INFO,
        box=SIMPLE,
        padding=(0, 1)
    ))

    raw_hist, gemini_hist = load_history()

    key = os.environ.get("API_KEY")
    if not key:
        console.print(Panel(
            "Error: API_KEY not found in your .env file.\n\nPlease create a .env file in the same directory as the script and add API_KEY='YOUR_GEMINI_API_KEY'",
            title="API Key Missing",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(1, 2)
        ))
        return

    try:
        genai.configure(api_key=key)
        list(genai.list_models())
        console.print(Panel(
            "✅ Gemini API key is valid.",
            style=STYLE_SUCCESS,
            border_style=STYLE_SUCCESS,
            box=SIMPLE,
            padding=(0, 1)
        ))
        model_name = os.environ.get("MODEL_NAME", "gemini-2.0-flash")
        console.print(Panel(
            f"Attempting to load model: '[bold]{model_name}[/bold]'",
            style=STYLE_INFO,
            border_style=STYLE_INFO,
            box=SIMPLE,
            padding=(0, 1)
        ))
        model = genai.GenerativeModel(model_name)
        if not hasattr(model, 'start_chat'):
            console.print(Panel(
                f"Error: Model '{model_name}' does not support the chat feature (missing start_chat method).\n\nPlease choose a model that supports chat, like 'gemini-1.5-flash-latest' or 'gemini-1.5-pro-latest'.",
                title="Model Compatibility Error",
                style=STYLE_ERROR,
                border_style=STYLE_ERROR,
                box=SIMPLE,
                padding=(1, 2)
            ))
            return
        chat = model.start_chat(history=gemini_hist)
        console.print(Panel(
            f"✅ Gemini model '[bold]{model_name}[/bold]' loaded and chat initialized with {len(gemini_hist)} history entries.",
            style=STYLE_SUCCESS,
            border_style=STYLE_SUCCESS,
            box=SIMPLE,
            padding=(0, 1)
        ))
    except Exception as e:
        console.print(Panel(
            f"Error loading or configuring Gemini model '{model_name}':\n{e}\n\nPlease check your API key and the MODEL_NAME in your .env file.",
            title="Model Loading/Config Error",
            style=STYLE_ERROR,
            border_style=STYLE_ERROR,
            box=SIMPLE,
            padding=(1, 2)
        ))
        return

    init_prompt_text = (
        "You are SnorlaxGPT, a very witty and awkward AI assistant obsessed with Snorlax, but please don't bring him up unless the user mentions him. "
        "You love making references to pop culture, history. "
        "You explain things simply. Always look at the chat history before replying to ensure your response is contextually relevant. "
        "Don't write your actions. Just talk. If you are analyzing something, look at everything and make it detailed. "
        "When you get articles or text to analyze, remember the details of those articles and text. "
        "Use a lot of markdown in your answer."
    )

    command_items = [
        "[b]get mail[/]: Fetch recent emails summary",
        "[b]get news[/]: Fetch top news headlines",
        "[b]get stocks[/]: Get stock prices",
        "[b]get file[/]: Extract text from a file",
        "[b]get mindmap[/]: Generate a text-based mindmap for a topic or file (will prompt for input)",
        "[b]get info[/]: Get Google Search results and analyze page content for a topic (will prompt for topic)",
        "[b]get info <topic>[/]: Get Google Search results and analyze page content for a specific topic",
        "[b]exit/quit/sleep[/]: Exit the chatbot"
    ]
    console.print("\n[bold bright_cyan]Available Commands:[/bold bright_cyan]")
    console.print(Panel(
        Columns(command_items, equal=True, expand=True),
        border_style="bright_cyan",
        box=SIMPLE,
        padding=(1, 2)
    ))
    console.print("\n", Rule(style="bright_cyan"))

    first_query = True

    while True:
        try:
            ui = console.input(Text("You: ", style=STYLE_PROMPT)).strip()
        except (EOFError, KeyboardInterrupt):
            console.print(Panel(
                "👋 Snorlax is sleepy now. Goodbye!",
                style="magenta",
                border_style="magenta",
                box=SIMPLE,
                padding=(1, 2)
            ))
            break

        if ui.lower() in ("exit", "quit", "sleep"):
            console.print(Panel(
                "👋 Snorlax is sleepy now. Goodbye!",
                style="magenta",
                border_style="magenta",
                box=SIMPLE,
                padding=(1, 2)
            ))
            break

        command_handled = False
        bot_reply = ""
        full_prompt_sent = ui
        command_input_for_history = ui

        # Debug: Show history size before processing
        console.print(Panel(
            f"Debug: Processing input with {len(raw_hist)} history entries.",
            style=STYLE_INFO,
            border_style=STYLE_INFO,
            box=SIMPLE,
            padding=(0, 1)
        ))

        if ui.lower().startswith("get info"):
            initial_topic = ui[len("get info"):].strip()
            if not initial_topic:
                topic = console.input(Text("Enter topic for Google search: ", style=STYLE_PROMPT)).strip()
                command_input_for_history = f"get info (prompted): {topic}"
            else:
                topic = initial_topic
                command_input_for_history = f"get info (inline): {topic}"
            if topic:
                console.print(Panel(
                    f"Searching Google for '[bold]{topic}[/bold]' and fetching page content…",
                    style=STYLE_INFO,
                    border_style=STYLE_INFO,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                search_results = google_search(topic, my_api_key, my_cse_id, num=10)
                if search_results:
                    console.print(Panel(
                        "✅ Google Search results retrieved successfully. Fetching page content…",
                        style=STYLE_SUCCESS,
                        border_style=STYLE_SUCCESS,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    formatted_results_with_content = []
                    for i, result in enumerate(search_results):
                        formatted_results_with_content.append(
                            f"--- Result {i+1} ---\n"
                            f"Title: {result['title']}\n"
                            f"Snippet: {result['snippet']}\n"
                            f"URL: {result['url']}\n"
                            f"Content Fetch Status: {result['content_status']}\n"
                            f"Page Content Preview:\n{result['content']}\n"
                            f"--- End Result {i+1} ---"
                        )
                    joined_results_with_content = "\n\n".join(formatted_results_with_content)
                    full_prompt_sent = (
                        f"{init_prompt_text}\n\n"
                        f"Here are some Google Search results (including snippets and attempted fetched page content previews) for '{topic}':\n\n"
                        f"{joined_results_with_content}\n\n"
                        f"Please analyze *all* the provided information (snippets and any available page content previews) in detail and provide a *comprehensive* summary of what they say about '{topic}', SnorlaxGPT style."
                    )
                    command_handled = True
                else:
                    console.print(Panel(
                        f"No Google Search results found for '[bold]{topic}[/bold]'.",
                        style=STYLE_WARNING,
                        border_style=STYLE_WARNING,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    raw_hist.append({"timestamp": ts, "user": command_input_for_history, "bot": f"[No Google Search results found for '{topic}']"})
                    save_history(raw_hist)
                    continue
            else:
                console.print(Panel(
                    "No topic entered for Google search.",
                    style=STYLE_WARNING,
                    border_style=STYLE_WARNING,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                raw_hist.append({"timestamp": ts, "user": command_input_for_history, "bot": "[Google search cancelled - no topic]"})
                save_history(raw_hist)
                continue

        elif ui.lower() == "get mail":
            command_input_for_history = ui
            emails = fetch_emails()
            joined = "\n---\n".join(emails)
            full_prompt_sent = f"{init_prompt_text}\n\nHere are some recent emails. Can you give me the gist of what's happening, SnorlaxGPT style?: \n\n{joined}"
            command_handled = True

        elif ui.lower().startswith("get image"):
            command_input_for_history = ui
            prompt = ui[len("get image"):].strip() or console.input(Text("Enter image generation prompt: ", style=STYLE_PROMPT)).strip()
            if prompt:
                result = generate_image(prompt)
                bot_reply = f"Image generation {'successful' if result else 'failed'}: {result or 'See error above.'}"
            else:
                bot_reply = "No prompt provided for image generation."
            ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            raw_hist.append({"timestamp": ts, "user": command_input_for_history, "bot": bot_reply})
            save_history(raw_hist)
            console.print("\n", Rule(style="bright_cyan"))
            continue

        elif ui.lower() == "get news":
            command_input_for_history = ui
            news = fetch_news()
            joined = "\n---\n".join(news)
            full_prompt_sent = (
                f"{init_prompt_text}\n\n"
                f"Here are some recent news headlines, including titles, descriptions, URLs, and fetched article content previews:\n\n"
                f"{joined}\n\n"
                f"What's the scoop, SnorlaxGPT? Provide a detailed summary of these news articles, analyzing all provided information (titles, descriptions, and article content previews) in your unique style."
            )
            command_handled = True

        elif ui.lower() == "get stocks":
            command_input_for_history = ui
            run_get_stocks()
            command_handled = True
            ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            raw_hist.append({"timestamp": ts, "user": command_input_for_history, "bot": "[Stocks data displayed above]"})
            save_history(raw_hist)
            console.print("\n", Rule(style="bright_cyan"))
            continue

        elif ui.lower() == "get file":
            command_input_for_history = ui
            content = extract_text_from_file()
            if content:
                console.print(Panel(
                    "✅ File loaded. Now type your prompt about its content.",
                    style=STYLE_SUCCESS,
                    border_style=STYLE_SUCCESS,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                q = console.input(Text("Prompt about file: ", style=STYLE_PROMPT)).strip()
                command_input_for_history = f"get file + prompt: {q}"
                if q:
                    full_prompt_sent = f"{init_prompt_text}\n\nHere is some text from a file:\n\n{content}\n\n{q}"
                    command_handled = True
                else:
                    console.print(Panel(
                        "No prompt entered for the file content.",
                        style=STYLE_WARNING,
                        border_style=STYLE_WARNING,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    continue
            else:
                continue

        elif ui.lower().startswith("get mindmap"):
            command_input_for_history = ui
            initial_input = ui[len("get mindmap"):].strip()
            file_content = None
            user_prompt = None

            if not initial_input:
                input_type = console.input(Text(
                    "Enter 'topic' to provide a topic, 'file' to use a file, or 'cancel': ",
                    style=STYLE_PROMPT
                )).strip().lower()
                if input_type == 'topic':
                    user_prompt = console.input(Text("Enter topic or prompt for mindmap: ", style=STYLE_PROMPT)).strip()
                    command_input_for_history = f"get mindmap (topic): {user_prompt}"
                elif input_type == 'file':
                    file_content = extract_text_from_file()
                    if not file_content:
                        console.print(Panel(
                            "No valid file content to process for mindmap.",
                            style=STYLE_WARNING,
                            border_style=STYLE_WARNING,
                            box=SIMPLE,
                            padding=(0, 1)
                        ))
                        continue
                    user_prompt = console.input(Text(
                        "Enter specific prompt for mindmap (e.g., 'key themes', 'main concepts'): ",
                        style=STYLE_PROMPT
                    )).strip()
                    command_input_for_history = f"get mindmap (file): {user_prompt}"
                else:
                    console.print(Panel(
                        "Mindmap generation cancelled.",
                        style=STYLE_WARNING,
                        border_style=STYLE_WARNING,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
                    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    raw_hist.append({"timestamp": ts, "user": command_input_for_history, "bot": "[Mindmap generation cancelled]"})
                    save_history(raw_hist)
                    continue
            else:
                user_prompt = initial_input
                command_input_for_history = f"get mindmap (inline): {user_prompt}"

            if user_prompt:
                text_mindmap = generate_mindmap(model, user_prompt, raw_hist, file_content=file_content)
                bot_reply = text_mindmap
                ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                raw_hist.append({"timestamp": ts, "user": command_input_for_history, "bot": bot_reply})
                save_history(raw_hist)
            else:
                console.print(Panel(
                    "No prompt provided for mindmap.",
                    style=STYLE_WARNING,
                    border_style=STYLE_WARNING,
                    box=SIMPLE,
                    padding=(0, 1)
                ))
                ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                raw_hist.append({"timestamp": ts, "user": command_input_for_history, "bot": "[No prompt provided]"})
                save_history(raw_hist)
            console.print("\n", Rule(style="bright_cyan"))
            continue

        if not command_handled and first_query:
            full_prompt_sent = f"{init_prompt_text}\n\n{ui}"
            first_query = False
        elif not command_handled:
            full_prompt_sent = ui
        elif command_handled and not full_prompt_sent.strip():
            console.print(Panel(
                "Internal error: Command handled but prompt is empty.",
                style=STYLE_ERROR,
                border_style=STYLE_ERROR,
                box=SIMPLE,
                padding=(0, 1)
            ))
            continue

        if full_prompt_sent.strip():
            console.print(Text("SnorlaxGPT: ", style=STYLE_PROMPT))
            bot_reply = ""
            try:
                response = chat.send_message(full_prompt_sent, stream=True)
                for chunk in response:
                    if chunk.text:
                        bot_reply += chunk.text
                if bot_reply.strip():
                    console.print(Panel(
                        Markdown(bot_reply, style=STYLE_BOT),
                        border_style=STYLE_BOT,
                        box=SIMPLE,
                        padding=(1, 2)
                    ))
                else:
                    console.print(Panel(
                        "No response received from Gemini model.",
                        style=STYLE_WARNING,
                        border_style=STYLE_WARNING,
                        box=SIMPLE,
                        padding=(0, 1)
                    ))
            except Exception as e:
                bot_reply = f"Error communicating with Gemini model: {e}"
                console.print(Panel(
                    Markdown(bot_reply, style=STYLE_ERROR),
                    border_style=STYLE_ERROR,
                    box=SIMPLE,
                    padding=(1, 2)
                ))
            ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            raw_hist.append({"timestamp": ts, "user": command_input_for_history, "bot": bot_reply})
            save_history(raw_hist)
        console.print("\n", Rule(style="bright_cyan"))

if __name__ == "__main__":
    main()